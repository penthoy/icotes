"""
PowerPoint document handler

Supports reading and writing PowerPoint files (.pptx)
"""

import io
import logging
from typing import Any, Dict, List, Optional, Union

from .base import (
    DocumentFormat,
    DocumentHandler,
    HandlerResult,
    register_handler,
)

logger = logging.getLogger(__name__)


class PPTXHandler(DocumentHandler):
    """Handler for PowerPoint documents (.pptx)"""
    
    @property
    def supported_formats(self) -> List[DocumentFormat]:
        return [DocumentFormat.POWERPOINT]
    
    async def read_content(
        self,
        file_data: bytes,
        file_path: str,
        options: Optional[Dict[str, Any]] = None
    ) -> HandlerResult:
        """
        Extract content from PowerPoint file.
        
        Options:
            max_slides: Maximum slides to read (default: 100)
            include_notes: Include speaker notes (default: True)
            include_tables: Extract tables separately (default: True)
        """
        options = options or {}
        max_slides = options.get("max_slides", 100)
        include_notes = options.get("include_notes", True)
        include_tables = options.get("include_tables", True)
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return HandlerResult(
                success=False,
                error="python-pptx not installed. Install with: pip install python-pptx"
            )
        
        try:
            pptx_file = io.BytesIO(file_data)
            prs = Presentation(pptx_file)
            
            content_parts = []
            tables = []
            total_slides = len(prs.slides)
            slides_to_read = min(total_slides, max_slides)
            
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx >= max_slides:
                    break
                
                content_parts.append(f"\n=== Slide {slide_idx + 1} ===\n")
                
                # Extract text from all shapes
                slide_text = []
                for shape in slide.shapes:
                    # Handle text frames
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Handle tables
                    if include_tables and shape.has_table:
                        table_data = self._extract_table(shape.table, slide_idx)
                        if table_data:
                            tables.append(table_data)
                            slide_text.append(f"[TABLE: {table_data['row_count']} rows]")
                
                content_parts.extend(slide_text)
                
                # Extract speaker notes
                if include_notes and slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text:
                        content_parts.append(f"\n[Notes: {notes_frame.text}]")
            
            # Truncation notice
            if total_slides > max_slides:
                content_parts.append(
                    f"\n... (showing {max_slides} of {total_slides} slides)"
                )
            
            # Extract metadata
            metadata = self._extract_metadata(prs)
            
            return HandlerResult(
                success=True,
                content="\n".join(content_parts),
                metadata=metadata,
                tables=tables,
                slides=total_slides
            )
            
        except Exception as e:
            logger.error(f"PowerPoint read error: {e}")
            return HandlerResult(
                success=False,
                error=f"Failed to read PowerPoint: {str(e)}"
            )
    
    def _extract_table(self, table, slide_idx: int) -> Optional[Dict[str, Any]]:
        """Extract data from a PowerPoint table"""
        try:
            rows_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip() if cell.text else "")
                rows_data.append(row_data)
            
            if not rows_data:
                return None
            
            # Use first row as headers
            headers = rows_data[0] if rows_data else []
            headers = [h if h else f"col_{i}" for i, h in enumerate(headers)]
            
            data = [
                dict(zip(headers, row))
                for row in rows_data[1:]
            ]
            
            return {
                "slide": slide_idx + 1,
                "columns": headers,
                "row_count": len(data),
                "data": data
            }
        except Exception as e:
            logger.debug(f"Table extraction error: {e}")
            return None
    
    def _extract_metadata(self, prs) -> Dict[str, Any]:
        """Extract presentation metadata"""
        metadata = {}
        
        try:
            core_props = prs.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "last_modified_by": core_props.last_modified_by or "",
            }
            
            # Add presentation dimensions
            metadata["slide_width"] = prs.slide_width
            metadata["slide_height"] = prs.slide_height
            
        except Exception as e:
            logger.debug(f"Metadata extraction error: {e}")
        
        return metadata
    
    async def write_content(
        self,
        content: Union[str, Dict[str, Any]],
        file_path: str,
        options: Optional[Dict[str, Any]] = None
    ) -> bytes:
        """
        Create PowerPoint from content.
        
        Content format:
            {
                "slides": [
                    {
                        "title": "Slide Title",
                        "content": ["Bullet 1", "Bullet 2"],
                        "notes": "Speaker notes"
                    },
                    ...
                ]
            }
        
        Options:
            title: Presentation title (for first slide)
            author: Author metadata
        """
        options = options or {}
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError(
                "python-pptx required for PowerPoint creation. "
                "Install with: pip install python-pptx"
            )
        
        prs = Presentation()
        
        # Set metadata
        prs.core_properties.title = options.get("title", "")
        prs.core_properties.author = options.get("author", "")
        
        if isinstance(content, str):
            # Simple text: create one slide per paragraph
            paragraphs = content.split("\n\n")
            for para in paragraphs:
                if para.strip():
                    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
                    title_shape = slide.shapes.title
                    body_shape = slide.placeholders[1]
                    
                    # First line as title
                    lines = para.strip().split("\n")
                    if lines:
                        title_shape.text = lines[0].lstrip("#").strip()
                        if len(lines) > 1:
                            tf = body_shape.text_frame
                            tf.text = lines[1]
                            for line in lines[2:]:
                                p = tf.add_paragraph()
                                p.text = line.strip()
        
        elif isinstance(content, dict):
            slides_data = content.get("slides", [])
            
            for slide_data in slides_data:
                if isinstance(slide_data, str):
                    # Simple text slide
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = slide_data[:50]
                    slide.placeholders[1].text = slide_data
                else:
                    # Structured slide
                    layout = prs.slide_layouts[1]  # Title and Content
                    slide = prs.slides.add_slide(layout)
                    
                    # Set title
                    if slide_data.get("title"):
                        slide.shapes.title.text = slide_data["title"]
                    
                    # Set content
                    body = slide.placeholders[1]
                    content_items = slide_data.get("content", [])
                    if content_items:
                        tf = body.text_frame
                        tf.text = content_items[0] if content_items else ""
                        for item in content_items[1:]:
                            p = tf.add_paragraph()
                            p.text = item
                    
                    # Add notes
                    if slide_data.get("notes"):
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = slide_data["notes"]
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()


# Register handler when module is imported
_handler = PPTXHandler()
register_handler(_handler)
